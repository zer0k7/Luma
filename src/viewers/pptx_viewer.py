"""Presentation (.pptx) viewer widget.

Parses OpenXML PresentationML documents using python-pptx and presents
slide-by-slide navigation with thumbnails and shape content.
"""

from pathlib import Path
from typing import Any, List

import gi
from pptx import Presentation  # type: ignore

gi.require_version("Gtk", "4.0")
from gi.repository import Gtk  # type: ignore # noqa: E402

from src.strings import (  # noqa: E402
    ERROR_PARSING_FAILED,
    PPTX_NEXT_SLIDE_TOOLTIP,
    PPTX_PREV_SLIDE_TOOLTIP,
    PPTX_SLIDE_STATUS_TEMPLATE,
)
from src.viewers.base import FormatViewerError  # noqa: E402


class PptxViewer(Gtk.Box):
    """Slide-by-slide presentation viewer with navigation controls."""

    def __init__(self, file_path: str) -> None:
        """Initialize PptxViewer with parsed slides.

        Args:
            file_path: Absolute filesystem path to .pptx presentation.

        Raises:
            FormatViewerError: If PPTX parsing fails.
        """
        super().__init__(orientation=Gtk.Orientation.VERTICAL, spacing=0)
        self.file_path = file_path
        self.current_index = 0
        self.slide_data: List[List[str]] = []

        self._load_presentation(file_path)
        self._build_toolbar()
        self._build_content_area()
        self._show_slide(0)

    def _load_presentation(self, file_path: str) -> None:
        """Extract slide texts from presentation file.

        Args:
            file_path: Absolute path to .pptx presentation.

        Raises:
            FormatViewerError: If pptx file is invalid or unreadable.
        """
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                slide_texts: List[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_texts.append(text)
                self.slide_data.append(slide_texts)
        except Exception as exc:
            filename = Path(file_path).name
            error_msg = ERROR_PARSING_FAILED.format(
                path=filename,
                format_name="PowerPoint Presentation (.pptx)",
            )
            raise FormatViewerError(error_msg, path=file_path) from exc

        if not self.slide_data:
            self.slide_data.append([])

    def _build_toolbar(self) -> None:
        """Construct top toolbar with slide navigation buttons."""
        self.toolbar = Gtk.Box(orientation=Gtk.Orientation.HORIZONTAL, spacing=6)
        self.toolbar.set_margin_top(6)
        self.toolbar.set_margin_bottom(6)
        self.toolbar.set_margin_start(12)
        self.toolbar.set_margin_end(12)

        self.btn_prev = Gtk.Button.new_from_icon_name("go-previous-symbolic")
        self.btn_prev.set_tooltip_text(PPTX_PREV_SLIDE_TOOLTIP)
        self.btn_prev.connect("clicked", self._on_prev_slide)
        self.toolbar.append(self.btn_prev)

        self.btn_next = Gtk.Button.new_from_icon_name("go-next-symbolic")
        self.btn_next.set_tooltip_text(PPTX_NEXT_SLIDE_TOOLTIP)
        self.btn_next.connect("clicked", self._on_next_slide)
        self.toolbar.append(self.btn_next)

        total = max(len(self.slide_data), 1)
        self.lbl_status = Gtk.Label(label=PPTX_SLIDE_STATUS_TEMPLATE.format(current=1, total=total))
        self.toolbar.append(self.lbl_status)

        self.append(self.toolbar)

    def _build_content_area(self) -> None:
        """Build the main slide canvas displaying the slide content."""
        self.scrolled = Gtk.ScrolledWindow()
        self.scrolled.set_vexpand(True)
        self.scrolled.set_hexpand(True)

        self.slide_box = Gtk.Box(orientation=Gtk.Orientation.VERTICAL, spacing=16)
        self.slide_box.set_margin_top(32)
        self.slide_box.set_margin_bottom(32)
        self.slide_box.set_margin_start(48)
        self.slide_box.set_margin_end(48)

        self.scrolled.set_child(self.slide_box)
        self.append(self.scrolled)

    def _show_slide(self, index: int) -> None:
        """Render the slide elements at the given index.

        Args:
            index: 0-based slide index.
        """
        # Clear previous slide children
        while child := self.slide_box.get_first_child():
            self.slide_box.remove(child)

        texts = self.slide_data[index] if index < len(self.slide_data) else []
        for i, text in enumerate(texts):
            lbl = Gtk.Label(label=text)
            lbl.set_wrap(True)
            lbl.set_xalign(0.0)
            if i == 0:
                lbl.add_css_class("title-1")
            else:
                lbl.add_css_class("body")
            self.slide_box.append(lbl)

        total = max(len(self.slide_data), 1)
        self.lbl_status.set_label(PPTX_SLIDE_STATUS_TEMPLATE.format(current=index + 1, total=total))

    def _on_prev_slide(self, _button: Any) -> None:
        """Go to the previous slide."""
        if self.current_index > 0:
            self.current_index -= 1
            self._show_slide(self.current_index)

    def _on_next_slide(self, _button: Any) -> None:
        """Go to the next slide."""
        if self.current_index < len(self.slide_data) - 1:
            self.current_index += 1
            self._show_slide(self.current_index)
